import linecache

from pptx import Presentation

from pptx.util import Inches, Pt


def set_presentation(filepath):
    return Presentation(filepath)


def set_home(presentation, title, author, title_rgb):
    prs = presentation
    oneSlide = prs.slides.add_slide(prs.slide_layouts[1])
    body_shapes = oneSlide.shapes.placeholders
    title_shape = body_shapes[0]
    title_frame = title_shape.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.font.color.rgb = title_rgb
    title_run.text = title



    author_shape = body_shapes[10]
    author_frame = author_shape.text_frame
    author_paragraph = author_frame.paragraphs[0]
    author_paragraph.text = author

    author_shape = body_shapes[10]
    author_frame = author_shape.text_frame
    author_paragraph = author_frame.paragraphs[0]
    author_paragraph.text = author


    return prs


def set_end(presentation, date_str, end_tittle_content):
    prs = presentation
    oneSlide = prs.slides.add_slide(prs.slide_layouts[3])
    body_shapes = oneSlide.shapes.placeholders
    title_shape = body_shapes[0]
    title_frame = title_shape.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = end_tittle_content

    subtitle_shape = body_shapes[13]
    subtitle_frame = subtitle_shape.text_frame
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.text = date_str
    return prs


def set_content(presentation, content_titles, current_i, content_font, current_content, other_content, current_rgb,
                other_rgb, content_number_rgb):
    prs = presentation
    content_back = prs.slide_layouts[2]
    slide = prs.slides.add_slide(content_back)
    one_col_left = 4.55
    two_col_left1 = 2.75
    two_col_left2 = 7.25
    width = 3.5
    height = 0.5
    tlen = len(content_titles)

    if tlen <= 3:
        height_total = height * tlen + 0.5 * (tlen - 1)
    else:
        height_total = height * (tlen / 2) + 0.5 * (tlen / 2 - 1)

    top = (7.5 - height_total) / 2
    top_flag = top
    for i in range(tlen):
        if i == 0 or (tlen >= 3 and i == int((tlen + 1) / 2)):
            i_top = top
        else:
            i_top = top_flag + (0.5 + height)

        top_flag = i_top

        i_left = 0
        if tlen <= 2:
            i_left = one_col_left
        elif tlen > 2 and i < tlen / 2:
            i_left = two_col_left1
        elif tlen > 2 and i >= tlen / 2:
            i_left = two_col_left2

        txBox = slide.shapes.add_textbox(Inches(i_left + 1), Inches(i_top), Inches(width), Inches(height))
        tf = txBox.text_frame
        title = tf.paragraphs[0]
        run = title.add_run()
        run.text = content_titles[i]

        font = run.font
        font.name = content_font
        font.bold = True
        img_path = other_content
        if i == current_i:
            font.color.rgb = current_rgb
            img_path = current_content
        else:
            font.color.rgb = other_rgb
        font.size = Pt(26)

        slide.shapes.add_picture(img_path, Inches(i_left), Inches(i_top))

        numberBox = slide.shapes.add_textbox(Inches(i_left + 0.125), Inches(i_top + 0.125), Inches(0.5), Inches(0.5))
        number = numberBox.text_frame
        number_title = number.paragraphs[0]
        number_run = number_title.add_run()
        number_run.text = str(i + 1)

        number_font = number_run.font
        number_font.name = content_font
        number_font.bold = True
        number_font.color.rgb = content_number_rgb
        number_font.size = Pt(18)

    return prs


def set_for_content(presentation, first_list, second_dict, content_list, second_title_rgb, content_font,
                    current_content,
                    other_content, current_rgb, other_rgb, content_number_rgb, text_img_path, import_font_size,
                    unsort_img_path, list_font_size):
    prs = presentation
    for title_i in range(len(first_list)):
        prs = set_content(presentation, first_list, title_i, content_font, current_content, other_content, current_rgb,
                          other_rgb, content_number_rgb)
        for second_title in second_dict[first_list[title_i]]:
            prs = gen_page(prs, second_title, content_list[second_title], second_title_rgb, text_img_path, content_font,
                           import_font_size, unsort_img_path, list_font_size)
    return prs


header = {}


def explain_markdown(file_path):
    first_list = []
    second_dict = {}

    content_list = {}

    txtfile = linecache.getlines(file_path)

    line_range = iter(range(len(txtfile)))
    is_code = False
    first_title = ""
    second_title = ""

    is_header = False

    for line_index in line_range:
        line = txtfile[line_index]
        line_strs = line.split(" ")
        if is_code == False:
            line = line.replace("\r", "")
            line = line.replace("\n", "")
            line = line.replace("\t", "")
        if len(line) == 0:
            continue
        elif line == "---":
            is_header = 1 - is_header
        elif is_header and line_strs[0] == "Title:":
            header['Title'] = line_strs[1].replace("\n", "")
        elif is_header and line_strs[0] == "Author:":
            header['Author'] = line_strs[1].replace("\n", "")
        elif is_header and line_strs[0] == "Date:":
            header['Date'] = line_strs[1].replace("\n", "")
        elif line_strs[0] == '#':
            first_title = line[2:]
            first_list.append(first_title)
            second_dict[first_title] = []
        elif line_strs[0] == '##':
            second_title = line[3:]
            second_dict[first_title].append(second_title)
            content_list[second_title] = []
        elif line_strs[0] == '-':
            content_list[second_title].append({
                "type": "list",
                "detail": line[2:]
            })
        elif line[:3] == "```":
            is_code = 1 - is_code
        elif is_code:
            continue
        elif line[:2] == '![':
            img_path_str = line.split('(')[-1].split(")")[0]
            content_list[second_title].append({
                "type": "image",
                "detail": img_path_str
            })
        elif line[:2] == '> ':
            content_list[second_title].append({
                "type": "text",
                "detail": line[2:]
            })
        else:
            continue

    return first_list, second_dict, content_list, header


def gen_page(prs, title, content_list, second_title_rgb, text_img_path, content_font, import_font_size, unsort_img_path,
             list_font_size):
    oneSlide = prs.slides.add_slide(prs.slide_layouts[0])
    body_shapes = oneSlide.shapes.placeholders
    title_shape = body_shapes[0]

    title_frame = title_shape.text_frame
    title_paragraph = title_frame.paragraphs[0]
    run = title_paragraph.add_run()
    run.text = title
    font = run.font
    font.color.rgb = second_title_rgb
    font.size = Pt(30)

    top_tag = 2
    is_have_img = False
    for content in content_list:
        if content['type'] == "text":
            oneSlide.shapes.add_picture(text_img_path, Inches(0.25), Inches(1.0))

            txBox = oneSlide.shapes.add_textbox(Inches(0.35), Inches(1.2), Inches(13), Inches(1))
            tf = txBox.text_frame
            title = tf.paragraphs[0]
            run = title.add_run()
            run.text = content['detail']
            font = run.font
            font.name = content_font
            font.size = Pt(import_font_size)
        if content['type'] == "list":
            if len(content_list) >= 3 and is_have_img == True:
                oneSlide.shapes.add_picture(unsort_img_path, Inches(7), Inches(top_tag))
                txBox = oneSlide.shapes.add_textbox(Inches(7.5), Inches(top_tag), Inches(5.5), Inches(0.25))
            else:
                oneSlide.shapes.add_picture(unsort_img_path, Inches(1), Inches(top_tag))
                txBox = oneSlide.shapes.add_textbox(Inches(1.5), Inches(top_tag), Inches(10), Inches(0.25))
            tf = txBox.text_frame
            title = tf.paragraphs[0]
            run = title.add_run()
            run.text = content['detail']
            font = run.font
            font.name = content_font
            font.size = Pt(list_font_size)
            font.bold = True
            top_tag = top_tag + 0.5
        if content['type'] == "image":
            detail_img_path = content['detail']
            if len(content_list) >= 3 and is_have_img == False:
                oneSlide.shapes.add_picture(detail_img_path, Inches(0.3), Inches(top_tag), width=Inches(6))
                is_have_img = True
            elif is_have_img == True:
                oneSlide.shapes.add_picture(detail_img_path, Inches(7), Inches(top_tag), width=Inches(6))
            else:
                oneSlide.shapes.add_picture(detail_img_path, Inches(1.5), Inches(top_tag), height=Inches(5))
    return prs


def gen_ppt_default(model_name, out_file_path, md_file_path, project_path, other_rgb, current_rgb, content_number_rgb,
                    second_title_rgb):
    # project_path = '..'
    # model_name = 'shuicai-model'
    # project_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    content_font = '微软雅黑'
    other_content = project_path + '/ppt-model/images/' + model_name + '/other_content.png'

    current_content = project_path + '/ppt-model/images/' + model_name + '/current_content.png'
    # other_rgb = RGBColor(126, 126, 126)
    # current_rgb = RGBColor(192, 125, 52)
    # content_number_rgb = RGBColor(255, 255, 255)
    # second_title_rgb = RGBColor(192, 125, 52)
    text_img_path = project_path + '/ppt-model/images/' + model_name + '/text_background.png'
    unsort_img_path = project_path + '/ppt-model/images/' + model_name + '/unsort_list.png'

    list_font_size = 18
    import_font_size = 16

    file_path = project_path + '/ppt-model/' + model_name + '.pptx'  # 模板的目录和文件名
    # out_file_path = '../out-print/readme.pptx'  # 输出的目录和文件名
    # md_file_path = '../md-file/readme/README.md'  # markdown的路径和文件名
    end_tittle_content = "感谢各位的聆听\n请领导批评指正"

    prs = set_presentation(file_path)
    first_list, second_dict, content_list, header = explain_markdown(md_file_path)

    prs = set_home(prs, header['Title'], header['Author'], current_rgb)
    prs = set_for_content(prs, first_list, second_dict, content_list, second_title_rgb, content_font, current_content,
                          other_content, current_rgb, other_rgb, content_number_rgb,
                          text_img_path, import_font_size, unsort_img_path, list_font_size)
    prs = set_end(prs, header['Date'], end_tittle_content)
    prs.save(out_file_path)

